import os
from flask import Flask, render_template, Response
import cv2
import numpy as np
import mediapipe as mp
from collections import deque
from pptx import Presentation

app = Flask(__name__)

class handDetector:
    def __init__(self, mode=False, maxHands=2, modelCom=1, detectionCon=0.5, trackingCon=0.5):
        self.mode = mode
        self.maxHands = maxHands
        self.modelCom = modelCom
        self.detectionCon = detectionCon
        self.trackingCon = trackingCon

        self.mpHands = mp.solutions.hands
        self.hands = self.mpHands.Hands(self.mode, self.maxHands, self.modelCom, self.detectionCon, self.trackingCon)
        self.mpDraw = mp.solutions.drawing_utils

    def findHands(self, img, draw=True):
        imgRGB = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
        self.results = self.hands.process(imgRGB)

        if self.results.multi_hand_landmarks:
            for handLms in self.results.multi_hand_landmarks:
                if draw:
                    self.mpDraw.draw_landmarks(img, handLms, self.mpHands.HAND_CONNECTIONS)
        return img

    def findPosition(self, img, handNo=0, draw=True):
        lmList = []
        if self.results.multi_hand_landmarks:
            myHand = self.results.multi_hand_landmarks[handNo]
            for id, lm in enumerate(myHand.landmark):
                h, w, c = img.shape
                cx, cy = int(lm.x * w), int(lm.y * h)
                lmList.append([id, cx, cy])
                if draw and id == 8:  # Only draw the index finger tip for clarity
                    cv2.circle(img, (cx, cy), 15, (255, 0, 255), cv2.FILLED)
        return lmList

def isHandOpen(lmList):
    if len(lmList) != 21:
        return False
    
    if lmList[8][2] > lmList[6][2]:
        return False
    if lmList[12][2] > lmList[10][2]:
        return False
    if lmList[16][2] > lmList[14][2]:
        return False
    if lmList[20][2] > lmList[18][2]:
        return False
    
    return True

def extract_powerpoint_slides(pptx_file):
    presentation = Presentation(pptx_file)
    slides = []

    for slide in presentation.slides:
        img = np.zeros((720, 1280, 3), np.uint8)  # Create a blank image
        shapes = slide.shapes

        for shape in shapes:
            if hasattr(shape, 'image'):
                image = shape.image
                image_bytes = image.blob
                nparr = np.frombuffer(image_bytes, np.uint8)
                img_np = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                img = cv2.resize(img_np, (1280, 720))  # Resize the image to fit the slide

        slides.append(img)

    return slides

@app.route('/meeting')
def index():
    return render_template('meeting.html')

global cap
cap = None

def gen():
    global cap
    pTime = 0
    cap = cv2.VideoCapture(0)
    detector = handDetector()
    pptx_file = r'presentations\Presentation1.pptx'  # Use the absolute path to the file

    if not os.path.exists(pptx_file):
        raise FileNotFoundError(f"{pptx_file} does not exist.")
        
    slides = extract_powerpoint_slides(pptx_file)
    slide_index = 0
    num_slides = len(slides)

    bpoints = [deque(maxlen=1024)]
    gpoints = [deque(maxlen=1024)]
    rpoints = [deque(maxlen=1024)]
    ypoints = [deque(maxlen=1024)]

    blue_index = 0
    green_index = 0
    red_index = 0
    yellow_index = 0

    colors = [(255, 0, 0), (0, 255, 0), (0, 0, 255), (0, 255, 255)]
    colorIndex = 0

    paintWindow = np.zeros((471, 636, 3)) + 255
    paintWindow = cv2.rectangle(paintWindow, (40, 1), (140, 65), (0, 0, 0), 2)
    paintWindow = cv2.rectangle(paintWindow, (160, 1), (255, 65), (255, 0, 0), 2)
    paintWindow = cv2.rectangle(paintWindow, (275, 1), (370, 65), (0, 255, 0), 2)
    paintWindow = cv2.rectangle(paintWindow, (390, 1), (485, 65), (0, 0, 255), 2)
    paintWindow = cv2.rectangle(paintWindow, (505, 1), (600, 65), (0, 255, 255), 2)

    cv2.putText(paintWindow, "CLEAR", (49, 33), cv2.FONT_HERSHEY_SIMPLEX, 0.5, (0, 0, 0), 2, cv2.LINE_AA)
    cv2.putText(paintWindow, "BLUE", (185, 33), cv2.FONT_HERSHEY_SIMPLEX, 0.5, (0, 0, 0), 2, cv2.LINE_AA)
    cv2.putText(paintWindow, "GREEN", (298, 33), cv2.FONT_HERSHEY_SIMPLEX, 0.5, (0, 0, 0), 2, cv2.LINE_AA)
    cv2.putText(paintWindow, "RED", (420, 33), cv2.FONT_HERSHEY_SIMPLEX, 0.5, (0, 0, 0), 2, cv2.LINE_AA)
    cv2.putText(paintWindow, "YELLOW", (520, 33), cv2.FONT_HERSHEY_SIMPLEX, 0.5, (0, 0, 0), 2, cv2.LINE_AA)

    while cap.isOpened():
        success, img = cap.read()
        if not success:
            break
        
        img = cv2.flip(img, 1)
        img = detector.findHands(img)
        lmList = detector.findPosition(img)

        if detector.results.multi_handedness:
            for idx, hand_handedness in enumerate(detector.results.multi_handedness):
                hand_label = hand_handedness.classification[0].label
                lmList = detector.findPosition(img, idx, draw=False)

                if len(lmList) != 0:
                    fore_finger = (lmList[8][1], lmList[8][2])
                    center = fore_finger
                    thumb = (lmList[4][1], lmList[4][2])

                    if hand_label == "Right":
                        cv2.circle(img, center, 3, (0, 255, 0), -1)

                        if (thumb[1] - center[1] < 30):
                            bpoints.append(deque(maxlen=512))
                            blue_index += 1
                            gpoints.append(deque(maxlen=512))
                            green_index += 1
                            rpoints.append(deque(maxlen=512))
                            red_index += 1
                            ypoints.append(deque(maxlen=512))
                            yellow_index += 1

                        elif center[1] <= 65:
                            if 40 <= center[0] <= 140:
                                bpoints = [deque(maxlen=512)]
                                gpoints = [deque(maxlen=512)]
                                rpoints = [deque(maxlen=512)]
                                ypoints = [deque(maxlen=512)]
                                blue_index = 0
                                green_index = 0
                                red_index = 0
                                yellow_index = 0
                                paintWindow[67:, :, :] = 255
                            elif 160 <= center[0] <= 255:
                                colorIndex = 0
                            elif 275 <= center[0] <= 370:
                                colorIndex = 1
                            elif 390 <= center[0] <= 485:
                                colorIndex = 2
                            elif 505 <= center[0] <= 600:
                                colorIndex = 3
                        else:
                            if colorIndex == 0 and blue_index < len(bpoints):
                                bpoints[blue_index].appendleft(center)
                            elif colorIndex == 1 and green_index < len(gpoints):
                                gpoints[green_index].appendleft(center)
                            elif colorIndex == 2 and red_index < len(rpoints):
                                rpoints[red_index].appendleft(center)
                            elif colorIndex == 3 and yellow_index < len(ypoints):
                                ypoints[yellow_index].appendleft(center)

                    elif hand_label == "Left":
                        if isHandOpen(lmList):
                            cx = (lmList[0][1] + lmList[5][1]) // 2
                            cy = (lmList[0][2] + lmList[5][2]) // 2
                            center = (cx, cy)   

                            # Paint a circle at the center of the left hand
                            cv2.circle(img, center, 0, (0, 255, 0), cv2.FILLED)

                            # Check if the center touches the specific region for clearing
                            if (thumb[1] - center[1]<30 ):
                                bpoints = [deque(maxlen=512)]
                                gpoints = [deque(maxlen=512)]
                                rpoints = [deque(maxlen=512)]
                                ypoints = [deque(maxlen=512)]
                                blue_index = 0
                                green_index = 0
                                red_index = 0
                                yellow_index = 0
                                paintWindow[67:, :, :] = 255

        points = [bpoints, gpoints, rpoints, ypoints]
        for i in range(len(points)):
            for j in range(len(points[i])):
                for k in range(1, len(points[i][j])):
                    if points[i][j][k - 1] is None or points[i][j][k] is None:
                        continue
                    cv2.line(paintWindow, points[i][j][k - 1], points[i][j][k], colors[i], 2)

        # Create a larger image that contains both the video feed and the whiteboard
        combined_img = np.zeros((max(img.shape[0], paintWindow.shape[0]), img.shape[1] + paintWindow.shape[1], 3), dtype=np.uint8)

        # Place the video feed and whiteboard side by side
        combined_img[:img.shape[0], :img.shape[1]] = img
        combined_img[:paintWindow.shape[0], img.shape[1]:] = paintWindow

        ret, jpeg = cv2.imencode('.jpg', combined_img)
        frame = jpeg.tobytes()
        yield (b'--frame\r\n'
               b'Content-Type: image/jpeg\r\n\r\n' + frame + b'\r\n\r\n')

@app.route('/video_feed')
def video_feed():
    return Response(gen(), mimetype='multipart/x-mixed-replace; boundary=frame')


@app.route('/')
def firstpage():
    return render_template('home.html')

@app.route('/home')
def home():
    return render_template('home.html')

@app.route('/simulation')
def simulation():
    return render_template('simulation.html')


@app.route('/service')
def service():
    return render_template('service.html')


@app.route('/stop_video')
def stop_video():
    global cap
    if cap is not None:
        cap.release()
    return render_template('home.html')

if __name__ == '__main__':
    app.run(debug=True)
